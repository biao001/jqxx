"""
build_ppt.py — 生成 DMS 行为识别算法 A 设计说明 PPT

5 页 16:9 蓝紫渐变风格，模仿"疲劳检测.pdf"。
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# ---------- 常量 ----------
BLUE        = RGBColor(0x4A, 0x7B, 0xF5)
BLUE_DARK   = RGBColor(0x2F, 0x5ACC, 0xDC) if False else RGBColor(0x2F, 0x5A, 0xCC)
PURPLE      = RGBColor(0x6B, 0x4F, 0xD9)
PURPLE_DARK = RGBColor(0x4A, 0x38, 0xB0)
LIGHT_BG    = RGBColor(0xF5, 0xF7, 0xFF)
CARD_BG     = RGBColor(0xFF, 0xFF, 0xFF)
CARD_BORDER = RGBColor(0xE5, 0xE7, 0xFF)
TEXT_DARK   = RGBColor(0x1F, 0x29, 0x37)
TEXT_MID    = RGBColor(0x4B, 0x55, 0x63)
TEXT_LIGHT  = RGBColor(0x6B, 0x72, 0x80)
ACCENT      = RGBColor(0x31, 0x45, 0xE0)

SLIDE_W = 13.333
SLIDE_H = 7.5
FONT_ZH = "Microsoft YaHei"
FONT_EN = "Segoe UI"


# ---------- 帮助函数 ----------

def set_fill(shape, rgb):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb


def set_no_line(shape):
    shape.line.fill.background()


def set_line(shape, rgb, width_pt=0.75):
    shape.line.color.rgb = rgb
    shape.line.width = Pt(width_pt)


def add_textbox(slide, x, y, w, h, text, size=14, bold=False, color=TEXT_DARK,
                align=PP_ALIGN.LEFT, font=FONT_ZH, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tb.text_frame.margin_left = Emu(0)
    tb.text_frame.margin_right = Emu(0)
    tb.text_frame.margin_top = Emu(0)
    tb.text_frame.margin_bottom = Emu(0)
    tb.text_frame.word_wrap = True
    tb.text_frame.vertical_anchor = anchor
    p = tb.text_frame.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return tb


def add_bullets(slide, x, y, w, h, items, size=11, color=TEXT_MID, spacing=4,
                font=FONT_ZH, bullet_color=PURPLE):
    """带自定义圆点的多行项，每行 small dot + 文字"""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(spacing)
        # 圆点
        r0 = p.add_run()
        r0.text = "● "
        r0.font.name = font
        r0.font.size = Pt(size)
        r0.font.color.rgb = bullet_color
        # 正文
        r1 = p.add_run()
        r1.text = it
        r1.font.name = font
        r1.font.size = Pt(size)
        r1.font.color.rgb = color
    return tb


def add_rounded_rect(slide, x, y, w, h, fill=CARD_BG, line=None, line_w=0.75,
                     shadow=False, radius=0.08):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(x), Inches(y), Inches(w), Inches(h))
    # 调整圆角半径
    try:
        shape.adjustments[0] = radius
    except Exception:
        pass
    set_fill(shape, fill)
    if line is None:
        set_no_line(shape)
    else:
        set_line(shape, line, line_w)
    if not shadow:
        _remove_shadow(shape)
    shape.text_frame.margin_left = Emu(0)
    shape.text_frame.margin_right = Emu(0)
    shape.text_frame.margin_top = Emu(0)
    shape.text_frame.margin_bottom = Emu(0)
    return shape


def _remove_shadow(shape):
    """移除默认黑色阴影"""
    sppr = shape._element.spPr
    # 移除已有 effectLst
    for el in sppr.findall(qn("a:effectLst")):
        sppr.remove(el)
    eff = etree.SubElement(sppr, qn("a:effectLst"))


def add_circle(slide, x, y, d, fill=BLUE, line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                               Inches(x), Inches(y), Inches(d), Inches(d))
    set_fill(s, fill)
    if line is None:
        set_no_line(s)
    else:
        set_line(s, line)
    _remove_shadow(s)
    return s


def add_number_badge(slide, x, y, w, h, num, fill=BLUE_DARK):
    """左上角编号小块，深蓝底白字"""
    s = add_rounded_rect(slide, x, y, w, h, fill=fill, radius=0.35)
    tf = s.text_frame
    tf.margin_left = Emu(30000); tf.margin_right = Emu(30000)
    tf.margin_top = Emu(10000);  tf.margin_bottom = Emu(10000)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.space_before = Pt(0); p.space_after = Pt(0)
    r = p.add_run(); r.text = num
    r.font.name = FONT_EN; r.font.size = Pt(12); r.font.bold = True
    r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    return s


def add_icon_circle(slide, cx, cy, d, text, fill=BLUE, text_color=None):
    """顶部圆图标 + Unicode 字符"""
    s = add_circle(slide, cx - d/2, cy - d/2, d, fill=fill)
    tf = s.text_frame
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0);  tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    p.space_before = Pt(0); p.space_after = Pt(0)
    r = p.add_run(); r.text = text
    r.font.name = FONT_EN; r.font.size = Pt(18); r.font.bold = True
    r.font.color.rgb = text_color or RGBColor(0xFF, 0xFF, 0xFF)
    return s


def add_arrow(slide, x, y, w, h, fill=BLUE):
    s = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                               Inches(x), Inches(y), Inches(w), Inches(h))
    set_fill(s, fill); set_no_line(s); _remove_shadow(s)
    return s


def add_title_bar(slide, title_text, subtitle_text):
    """左上 标题（蓝竖线 + 主标题 + 副标题 + 下方水平分隔线）"""
    # 蓝色粗竖线
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(0.55), Inches(0.45),
                                 Inches(0.12), Inches(0.6))
    set_fill(bar, BLUE); set_no_line(bar); _remove_shadow(bar)
    # 主标题
    add_textbox(slide, 0.78, 0.38, 7.0, 0.78, title_text,
                size=32, bold=True, color=TEXT_DARK, font=FONT_ZH)
    # 副标题
    add_textbox(slide, 0.78, 1.08, 10.0, 0.38, subtitle_text,
                size=14, color=TEXT_LIGHT, font=FONT_ZH)
    # 水平分隔线 + 右端蓝点
    ln = slide.shapes.add_connector(1, Inches(0.78), Inches(1.55),
                                    Inches(8.6), Inches(1.55))
    ln.line.color.rgb = RGBColor(0xD5, 0xD8, 0xE6)
    ln.line.width = Pt(1)
    # 蓝点作为线端
    add_circle(slide, 8.55, 1.48, 0.14, fill=BLUE)


def add_camera_decoration(slide):
    """右上角装饰: 紫色小方块阵 + 蓝色镜头圈"""
    # 像素方块阵（3x4 紫色小方块）
    for r in range(3):
        for c in range(4):
            x = 11.05 + c * 0.14
            y = 0.45 + r * 0.14
            sq = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                        Inches(x), Inches(y),
                                        Inches(0.10), Inches(0.10))
            shade = [PURPLE, BLUE, RGBColor(0xB7, 0xB0, 0xEE)][
                (r + c) % 3]
            set_fill(sq, shade); set_no_line(sq); _remove_shadow(sq)
    # 大镜头外圆（浅灰）
    outer = add_circle(slide, 11.85, 0.30, 0.95, fill=RGBColor(0xEE, 0xEF, 0xF8))
    # 外环
    ring = add_circle(slide, 11.92, 0.37, 0.80, fill=RGBColor(0x2D, 0x38, 0x7A))
    # 中心镜头
    lens = add_circle(slide, 12.02, 0.47, 0.60, fill=RGBColor(0x0F, 0x16, 0x40))
    # 高光
    hi = add_circle(slide, 12.15, 0.55, 0.18, fill=RGBColor(0x4A, 0x7B, 0xF5))


def add_bottom_bar(slide, chips):
    """底部胶囊状状态栏"""
    bar = add_rounded_rect(slide, 1.2, 6.55, 11.0, 0.7,
                           fill=RGBColor(0xFB, 0xFC, 0xFF),
                           line=RGBColor(0xD0, 0xD6, 0xF2), line_w=1.0,
                           radius=0.45)
    # 若干 chip, 等分
    n = len(chips)
    seg = 11.0 / n
    for i, (icon_char, text) in enumerate(chips):
        cx = 1.2 + seg * i + seg / 2
        # 图标小圆
        add_icon_circle(slide, cx - 1.5, 6.9, 0.42, icon_char,
                        fill=BLUE if i % 2 == 0 else PURPLE)
        # 文字
        add_textbox(slide, cx - 1.2, 6.72, 2.4, 0.4, text,
                    size=15, bold=True, color=TEXT_DARK,
                    align=PP_ALIGN.LEFT,
                    anchor=MSO_ANCHOR.MIDDLE)
        # 分隔线（除最后一个）
        if i < n - 1:
            x = 1.2 + seg * (i + 1)
            ln = slide.shapes.add_connector(1, Inches(x), Inches(6.75),
                                             Inches(x), Inches(7.05))
            ln.line.color.rgb = RGBColor(0xD5, 0xDA, 0xEE)
            ln.line.width = Pt(0.75)


# ---------- 卡片模板 ----------

def add_content_card(slide, x, y, w, h, num, title, bullets,
                     icon_char="★", icon_fill=BLUE):
    # 主卡片
    card = add_rounded_rect(slide, x, y + 0.5, w, h - 0.5,
                            fill=CARD_BG, line=CARD_BORDER, line_w=1.0,
                            radius=0.08)
    # 圆形图标位于卡片上沿（露出一半）
    add_icon_circle(slide, x + w / 2, y + 0.5, 0.8, icon_char, fill=icon_fill)
    # 编号徽章
    add_number_badge(slide, x + 0.22, y + 1.0, 0.55, 0.35, num,
                     fill=BLUE_DARK)
    # 标题
    add_textbox(slide, x + 0.22, y + 1.40, w - 0.3, 0.5, title,
                size=20, bold=True, color=TEXT_DARK, font=FONT_ZH,
                align=PP_ALIGN.LEFT)
    # 标题下划线（紫色短横）
    ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Inches(x + 0.22), Inches(y + 1.92),
                                Inches(0.6), Inches(0.04))
    set_fill(ln, PURPLE); set_no_line(ln); _remove_shadow(ln)
    # 项目列表
    add_bullets(slide, x + 0.22, y + 2.08, w - 0.4, h - 2.0, bullets,
                size=11, color=TEXT_MID, spacing=5)


# ========== Slide 1 : 项目定位 ==========

def build_slide_1(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    # 背景
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                            prs.slide_width, prs.slide_height)
    set_fill(bg, LIGHT_BG); set_no_line(bg); _remove_shadow(bg)

    add_title_bar(s, "项目定位", "基于摄像头视频流的本地危险驾驶行为识别模块")
    add_camera_decoration(s)

    # 4 张卡片
    cards = [
        ("01", "项目聚焦", "◎", BLUE, [
            "仅聚焦行为动作识别",
            "不涉及疲劳/视线模块",
            "面向本地部署与实验演示",
        ]),
        ("02", "项目目标", "▶", PURPLE, [
            "基于摄像头视频流实时检测",
            "输出手机/打电话/吸烟/安全带等事件",
            "形成风险评分与预警等级",
        ]),
        ("03", "关键检测对象", "◉", BLUE, [
            "手机使用 / 打电话",
            "吸烟",
            "未系安全带",
            "双手离开方向盘",
        ]),
        ("04", "核心特征指标", "✦", PURPLE, [
            "YOLOv8 bbox 置信度",
            "手机-头部距离 / 手腕位置",
            "方向盘 ROI 命中",
            "加权互补风险评分",
            "持续时间 / 时序滑窗",
        ]),
    ]
    x0 = 0.6; y0 = 2.0; card_w = 2.85; card_h = 4.3; gap = 0.25
    for i, (num, title, ic, ifill, bl) in enumerate(cards):
        x = x0 + i * (card_w + gap)
        add_content_card(s, x, y0, card_w, card_h, num, title, bl, ic, ifill)
        # 箭头（最后一个不画）
        if i < len(cards) - 1:
            ax = x + card_w + 0.02
            ay = y0 + card_h / 2 + 0.1
            add_arrow(s, ax, ay - 0.12, 0.22, 0.24,
                      fill=RGBColor(0xB5, 0xBF, 0xEB))

    add_bottom_bar(s, [("⏱", "实时监测"),
                       ("⬛", "本地部署"),
                       ("✓", "可解释评分")])


# ========== Slide 2 : 全流程技术路线 ==========

def build_slide_2(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                            prs.slide_width, prs.slide_height)
    set_fill(bg, LIGHT_BG); set_no_line(bg); _remove_shadow(bg)

    add_title_bar(s, "全流程技术路线", "围绕本地行为识别模块的完整处理链路")
    add_camera_decoration(s)

    cards = [
        ("01", "数据输入", "▤", BLUE, [
            "摄像头实时视频流",
            "Roboflow 公开数据集",
            "本地采集场景数据",
        ]),
        ("02", "数据合并", "⊞", PURPLE, [
            "3 数据集 → 8 类",
            "9137 张统一标注",
            "YOLO 格式 data.yaml",
        ]),
        ("03", "模型训练", "⚙", BLUE, [
            "YOLOv8n 基座",
            "AdamW + Mosaic/Mixup",
            "imgsz=640  epochs=40",
        ]),
        ("04", "实时推理", "▶", PURPLE, [
            "低光 CLAHE 预处理",
            "驾驶员 crop + 跳帧",
            "YOLO + Pose 并行",
        ]),
        ("05", "规则评分", "∑", BLUE, [
            "phone/calling 互斥",
            "方向盘 ROI 判定",
            "5 帧滑窗 + 加权互补",
        ]),
        ("06", "本地输出", "◱", PURPLE, [
            "监控风格 UI",
            "ACTIVE ALERTS 面板",
            "JSON 日志与告警",
        ]),
    ]
    x0 = 0.35; y0 = 2.0; card_w = 1.95; card_h = 4.3; gap = 0.13
    for i, (num, title, ic, ifill, bl) in enumerate(cards):
        x = x0 + i * (card_w + gap)
        add_content_card(s, x, y0, card_w, card_h, num, title, bl, ic, ifill)
        if i < len(cards) - 1:
            ax = x + card_w - 0.05
            ay = y0 + card_h / 2 + 0.15
            add_arrow(s, ax, ay - 0.11, 0.18, 0.22,
                      fill=RGBColor(0xB5, 0xBF, 0xEB))

    add_bottom_bar(s, [("▶", "视频流驱动"),
                       ("~", "时序分析"),
                       ("⚙", "本地实时推理")])


# ========== Slide 3 : 输入/输出设计 ==========

def build_slide_3(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                            prs.slide_width, prs.slide_height)
    set_fill(bg, LIGHT_BG); set_no_line(bg); _remove_shadow(bg)

    add_title_bar(s, "输入 / 输出设计", "面向行为识别模块的标准输入与分层输出")
    add_camera_decoration(s)

    # 左侧 "输入" 大容器 + 3 小卡
    in_container = add_rounded_rect(s, 0.55, 2.05, 3.6, 4.4,
                                    fill=RGBColor(0xFB, 0xFC, 0xFF),
                                    line=CARD_BORDER, line_w=1.0,
                                    radius=0.06)
    # 输入标签
    add_icon_circle(s, 0.9, 2.15, 0.5, "↓", fill=BLUE)
    add_textbox(s, 1.3, 1.98, 2.0, 0.4, "输入",
                size=18, bold=True, color=TEXT_DARK,
                anchor=MSO_ANCHOR.MIDDLE)

    in_cards = [
        ("✚", BLUE, "实时摄像头视频流",
         "本地实时采集  BGR uint8  ≥640×480"),
        ("≡", PURPLE, "Roboflow 公开数据集",
         "distracted_driving + cigarette + seatbelt  9137 张 8 类"),
        ("▤", BLUE, "本地采集数据",
         "光照/角度/戴眼镜等场景补充"),
    ]
    for i, (ic, col, t, sub) in enumerate(in_cards):
        y = 2.5 + i * 1.22
        add_rounded_rect(s, 0.75, y, 3.2, 1.0,
                         fill=CARD_BG, line=CARD_BORDER, line_w=0.8,
                         radius=0.1)
        add_icon_circle(s, 1.1, y + 0.5, 0.55, ic, fill=col)
        add_textbox(s, 1.55, y + 0.18, 2.4, 0.36, t,
                    size=13, bold=True, color=TEXT_DARK,
                    anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(s, 1.55, y + 0.55, 2.4, 0.35, sub,
                    size=9.5, color=TEXT_LIGHT,
                    anchor=MSO_ANCHOR.TOP)

    # 中间核心模块（圆形 + 方块）
    add_circle(s, 4.85, 3.55, 1.3, fill=PURPLE)
    add_icon_circle(s, 5.5, 3.9, 0.7, "⚙", fill=RGBColor(0x8A, 0x6E, 0xED))
    # 主核心方块
    core = add_rounded_rect(s, 4.35, 4.5, 2.3, 1.35,
                            fill=CARD_BG, line=PURPLE, line_w=1.2,
                            radius=0.08)
    add_textbox(s, 4.4, 4.6, 2.2, 0.4, "行为识别模块",
                size=14, bold=True, color=PURPLE_DARK,
                align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(s, 4.4, 5.0, 2.2, 0.8,
                "YOLOv8\n+ 规则判定\n+ 时序滑窗 + 评分决策",
                size=9, color=TEXT_MID,
                align=PP_ALIGN.CENTER)

    # 左右箭头装饰
    add_arrow(s, 4.05, 3.8, 0.28, 0.25, fill=RGBColor(0xB5, 0xBF, 0xEB))
    add_arrow(s, 6.95, 3.8, 0.28, 0.25, fill=RGBColor(0xB5, 0xBF, 0xEB))

    # 右侧 "输出" 容器 + 2x2 象限
    out_container = add_rounded_rect(s, 7.45, 2.05, 5.3, 4.4,
                                     fill=RGBColor(0xFB, 0xFC, 0xFF),
                                     line=CARD_BORDER, line_w=1.0,
                                     radius=0.06)
    add_icon_circle(s, 7.8, 2.15, 0.5, "↑", fill=PURPLE)
    add_textbox(s, 8.2, 1.98, 2.0, 0.4, "输出",
                size=18, bold=True, color=TEXT_DARK,
                anchor=MSO_ANCHOR.MIDDLE)

    quads = [
        ("A", BLUE, "原始检测", "•  bbox / class_id\n•  class_name\n•  confidence"),
        ("B", PURPLE, "行为事件",
         "•  phone_use / calling\n•  smoking / no_seatbelt\n•  hands_off_wheel"),
        ("C", BLUE, "风险量化",
         "•  Risk Score 0-100\n•  Tier: safe / warning\n     / danger / critical"),
        ("D", PURPLE, "展示输出",
         "•  监控界面 / 状态灯\n•  语音告警\n•  JSON 日志"),
    ]
    for i, (letter, col, title, body) in enumerate(quads):
        r, c = i // 2, i % 2
        qx = 7.65 + c * 2.5
        qy = 2.55 + r * 1.85
        add_rounded_rect(s, qx, qy, 2.4, 1.65,
                         fill=CARD_BG, line=CARD_BORDER, line_w=0.8,
                         radius=0.08)
        add_rounded_rect(s, qx + 0.15, qy + 0.15, 0.4, 0.4,
                         fill=col, radius=0.25)
        add_textbox(s, qx + 0.15, qy + 0.15, 0.4, 0.4, letter,
                    size=12, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                    font=FONT_EN)
        add_textbox(s, qx + 0.65, qy + 0.18, 1.8, 0.35, title,
                    size=13, bold=True, color=TEXT_DARK,
                    anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(s, qx + 0.15, qy + 0.65, 2.2, 1.0, body,
                    size=9.5, color=TEXT_MID)


# ========== Slide 4 : 技术架构 ==========

def build_slide_4(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                            prs.slide_width, prs.slide_height)
    set_fill(bg, LIGHT_BG); set_no_line(bg); _remove_shadow(bg)

    add_title_bar(s, "技术架构", "数据层—感知层—规则时序层—决策层的四层设计")
    add_camera_decoration(s)

    cards = [
        ("01", "数据层", "▤", BLUE, [
            "3 个 Roboflow 公开数据集",
            "9137 张 8 类统一标注",
            "实时视频流输入",
        ]),
        ("02", "感知层", "◎", PURPLE, [
            "YOLOv8n unified.pt",
            "YOLOv8n-pose 17 点辅助",
            "低光 CLAHE 预处理",
        ]),
        ("03", "规则时序层", "⏱", BLUE, [
            "calling ↔ phone_use 互斥",
            "方向盘 ROI 判定",
            "5 帧滑窗去抖 (3/3)",
            "duration_s 统计",
        ]),
        ("04", "决策层", "⚖", PURPLE, [
            "加权互补风险公式",
            "5 档风险等级映射",
            "高危 >2s 升级告警",
            "JSON 日志输出",
        ]),
    ]
    x0 = 0.55; y0 = 2.0; card_w = 2.35; card_h = 4.3; gap = 0.2
    for i, (num, title, ic, ifill, bl) in enumerate(cards):
        x = x0 + i * (card_w + gap)
        add_content_card(s, x, y0, card_w, card_h, num, title, bl, ic, ifill)
        if i < len(cards) - 1:
            ax = x + card_w - 0.03
            ay = y0 + card_h / 2 + 0.1
            add_arrow(s, ax, ay - 0.12, 0.22, 0.24,
                      fill=RGBColor(0xB5, 0xBF, 0xEB))

    # 右侧 关键能力
    cap_x = 10.8; cap_y = 2.0; cap_w = 2.1; cap_h = 4.3
    add_rounded_rect(s, cap_x, cap_y, cap_w, cap_h,
                     fill=RGBColor(0xFB, 0xFC, 0xFF),
                     line=CARD_BORDER, line_w=1.0, radius=0.08)
    # 顶部徽标
    add_rounded_rect(s, cap_x + 0.25, cap_y - 0.22, 1.6, 0.5,
                     fill=BLUE, radius=0.35)
    add_textbox(s, cap_x + 0.3, cap_y - 0.22, 1.5, 0.5,
                "◆ 关键能力", size=13, bold=True,
                color=RGBColor(0xFF, 0xFF, 0xFF),
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    caps = [
        ("⚡", "本地实时推理 25+ FPS"),
        ("≡", "可解释风险评分"),
        ("◉", "多行为并发检测"),
        ("⏱", "时序滑窗去抖"),
        ("➤", "支持后续扩展"),
    ]
    for i, (icn, txt) in enumerate(caps):
        y = cap_y + 0.55 + i * 0.73
        add_icon_circle(s, cap_x + 0.4, y + 0.22, 0.45, icn,
                        fill=PURPLE if i % 2 else BLUE)
        add_textbox(s, cap_x + 0.75, y, cap_w - 0.85, 0.5, txt,
                    size=10.5, bold=True, color=TEXT_DARK,
                    anchor=MSO_ANCHOR.MIDDLE)


# ========== Slide 5 : 效果呈现 ==========

def build_slide_5(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                            prs.slide_width, prs.slide_height)
    set_fill(bg, LIGHT_BG); set_no_line(bg); _remove_shadow(bg)

    add_title_bar(s, "效果呈现", "本地部署后的监控界面与结果展示")
    add_camera_decoration(s)

    # 中央 深色模拟 UI
    ui = add_rounded_rect(s, 3.35, 1.95, 6.7, 4.4,
                          fill=RGBColor(0x17, 0x1B, 0x2E),
                          radius=0.05)
    # 标题栏
    title_bar = add_rounded_rect(s, 3.35, 1.95, 6.7, 0.55,
                                 fill=RGBColor(0x0F, 0x14, 0x23),
                                 radius=0.05)
    add_textbox(s, 3.55, 1.95, 4.0, 0.55, "◆ DMS MONITOR",
                size=13, bold=True,
                color=RGBColor(0xBE, 0xCF, 0xFF),
                font=FONT_EN, anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(s, 7.5, 1.95, 2.4, 0.55, "● LIVE   10:24:35",
                size=10, color=RGBColor(0x80, 0xFF, 0x80),
                font=FONT_EN, anchor=MSO_ANCHOR.MIDDLE,
                align=PP_ALIGN.RIGHT)

    # 左上 视频框
    vid = add_rounded_rect(s, 3.55, 2.65, 3.3, 2.35,
                           fill=RGBColor(0x22, 0x2A, 0x45), radius=0.03)
    add_textbox(s, 3.65, 2.75, 1.5, 0.3, "● LIVE",
                size=9, bold=True, color=RGBColor(0x80, 0xFF, 0x80),
                font=FONT_EN)
    # 人脸占位圆
    add_circle(s, 4.7, 3.15, 1.0, fill=RGBColor(0x35, 0x3D, 0x5A))
    # 人体占位
    body = add_rounded_rect(s, 4.3, 3.95, 1.85, 1.0,
                            fill=RGBColor(0x3B, 0x41, 0x62), radius=0.2)
    # bbox 四角强调
    for (cx, cy) in [(4.0, 2.9), (5.45, 2.9), (4.0, 4.3), (5.45, 4.3)]:
        for dx, dy in [(0.15, 0), (0, 0.15)]:
            ln = s.shapes.add_connector(1,
                Inches(cx), Inches(cy),
                Inches(cx + dx), Inches(cy + dy))
            ln.line.color.rgb = RGBColor(0xFF, 0x64, 0x64)
            ln.line.width = Pt(2.0)
    # 标签
    add_rounded_rect(s, 4.0, 2.65, 1.0, 0.22,
                     fill=RGBColor(0xFF, 0x64, 0x64), radius=0.3)
    add_textbox(s, 4.0, 2.65, 1.0, 0.22, "calling",
                size=8, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
                font=FONT_EN, align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)

    # 右上 指标区
    metrics = [
        ("FPS", "28"),
        ("LATENCY", "35ms"),
        ("FRAME", "#1024"),
    ]
    for i, (k, v) in enumerate(metrics):
        y = 2.7 + i * 0.5
        add_textbox(s, 7.1, y, 1.5, 0.3, k,
                    size=9, color=RGBColor(0x9A, 0xAC, 0xD8), font=FONT_EN)
        add_textbox(s, 8.55, y, 1.4, 0.3, v,
                    size=12, bold=True,
                    color=RGBColor(0xBE, 0xCF, 0xFF), font=FONT_EN,
                    align=PP_ALIGN.RIGHT)

    # ACTIVE ALERTS 面板
    ap = add_rounded_rect(s, 7.0, 4.25, 3.0, 0.95,
                          fill=RGBColor(0x22, 0x2A, 0x45), radius=0.06)
    add_textbox(s, 7.1, 4.28, 2.0, 0.28, "▲ ACTIVE ALERTS",
                size=9, bold=True,
                color=RGBColor(0xFF, 0xC8, 0x64),
                font=FONT_EN)
    # 事件 1
    add_circle(s, 7.12, 4.6, 0.12, fill=RGBColor(0xFF, 0x64, 0x64))
    add_textbox(s, 7.3, 4.55, 2.6, 0.28,
                "calling   0.87   2.4s",
                size=9, color=RGBColor(0xDD, 0xE4, 0xFF), font=FONT_EN)
    # 事件 2
    add_circle(s, 7.12, 4.85, 0.12, fill=RGBColor(0xFF, 0xC8, 0x64))
    add_textbox(s, 7.3, 4.8, 2.6, 0.28,
                "no_seatbelt   0.75   3.1s",
                size=9, color=RGBColor(0xDD, 0xE4, 0xFF), font=FONT_EN)

    # 底部 RISK 条
    add_textbox(s, 3.55, 5.25, 2.0, 0.3, "RISK 68/100",
                size=11, bold=True, color=RGBColor(0xFF, 0xC8, 0x64),
                font=FONT_EN)
    risk_bg = add_rounded_rect(s, 3.55, 5.55, 4.5, 0.2,
                               fill=RGBColor(0x0B, 0x12, 0x20),
                               radius=0.4)
    risk_fill = add_rounded_rect(s, 3.55, 5.55, 3.05, 0.2,
                                 fill=RGBColor(0xFF, 0xC8, 0x64),
                                 radius=0.4)
    add_textbox(s, 8.15, 5.52, 1.8, 0.3, "[ WARNING ]",
                size=9, bold=True,
                color=RGBColor(0xFF, 0xC8, 0x64),
                font=FONT_EN)

    # 底部大字主标签
    add_textbox(s, 3.35, 5.85, 6.7, 0.45, "CALLING",
                size=24, bold=True,
                color=RGBColor(0xFF, 0x64, 0x64),
                font=FONT_EN, align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)
    # REC
    add_circle(s, 9.7, 6.05, 0.13, fill=RGBColor(0xFF, 0x64, 0x64))
    add_textbox(s, 9.85, 5.98, 0.5, 0.25, "REC",
                size=8, bold=True,
                color=RGBColor(0xFF, 0x64, 0x64),
                font=FONT_EN)

    # === 5 个箭头标注（左 2 + 右 3）===
    call_boxes = [
        # (x, y, title, body, 指向的箭头方向, color)
        (0.5, 2.05, "01", "实时视频 + bbox",
         "四角强调线，类别标签 calling", "right", BLUE),
        (0.5, 4.25, "02", "ACTIVE ALERTS 面板",
         "多行为并列，含置信度与持续时间", "right", PURPLE),
        (10.5, 2.05, "03", "风险仪表条 0-100",
         "5 档 tier 映射，颜色动态", "left", BLUE),
        (10.5, 3.70, "04", "状态指示灯",
         "DRIVER / CAMERA / SMOKE-DET", "left", PURPLE),
        (10.5, 5.35, "05", "日志与录制",
         "JSON 行式追加 + 截图", "left", BLUE),
    ]
    for (x, y, num, t, b, direction, col) in call_boxes:
        card = add_rounded_rect(s, x, y, 2.6, 1.3,
                                fill=CARD_BG, line=CARD_BORDER, line_w=0.8,
                                radius=0.12)
        # 编号
        add_rounded_rect(s, x + 0.15, y + 0.15, 0.4, 0.4,
                         fill=col, radius=0.4)
        add_textbox(s, x + 0.15, y + 0.15, 0.4, 0.4, num,
                    size=11, bold=True,
                    color=RGBColor(0xFF, 0xFF, 0xFF),
                    font=FONT_EN, align=PP_ALIGN.CENTER,
                    anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(s, x + 0.62, y + 0.15, 2.0, 0.4, t,
                    size=12, bold=True, color=TEXT_DARK,
                    anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(s, x + 0.15, y + 0.6, 2.35, 0.7, b,
                    size=9.5, color=TEXT_MID)

        # 小箭头指向中央 UI
        if direction == "right":
            add_arrow(s, x + 2.6, y + 0.45, 0.55, 0.18,
                      fill=RGBColor(0xB5, 0xBF, 0xEB))
        else:
            # 左指箭头：用 LEFT_ARROW
            ar = s.shapes.add_shape(MSO_SHAPE.LEFT_ARROW,
                                     Inches(x - 0.55), Inches(y + 0.45),
                                     Inches(0.55), Inches(0.18))
            set_fill(ar, RGBColor(0xB5, 0xBF, 0xEB))
            set_no_line(ar); _remove_shadow(ar)

    add_bottom_bar(s, [("⏱", "实时监测"),
                       ("⬛", "本地部署"),
                       ("☰", "实验分析")])


# ---------- 主入口 ----------

def main():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    build_slide_1(prs)
    build_slide_2(prs)
    build_slide_3(prs)
    build_slide_4(prs)
    build_slide_5(prs)

    out = r"D:\Desktop\DMS\behavior_algo_a\docs\行为识别算法A_设计说明.pptx"
    prs.save(out)

    import os
    size_kb = os.path.getsize(out) / 1024
    print(f"[OK] 保存: {out}  ({size_kb:.1f} KB)")


if __name__ == "__main__":
    main()
